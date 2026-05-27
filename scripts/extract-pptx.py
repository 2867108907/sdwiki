#!/usr/bin/env python3
"""
Extract text from PowerPoint pptx file for SDWiki ingestion
Usage: ./scripts/extract-pptx.py input.pptx
"""

import sys
import os

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx not installed", file=sys.stderr)
    print("Install with: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

def main():
    if len(sys.argv) != 2:
        print("Usage: ./extract-pptx.py <input.pptx>", file=sys.stderr)
        sys.exit(1)

    input_file = sys.argv[1]

    if not os.path.exists(input_file):
        print(f"Error: Input file {input_file} not found", file=sys.stderr)
        sys.exit(1)

    ext = os.path.splitext(input_file)[1].lower()
    if ext != '.pptx':
        print(f"Error: Only .pptx files are supported, got {ext}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(input_file)
    lines = []

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                lines.append(shape.text)

    print('\n'.join(lines))
    print(f"\n--- Extracted {len(prs.slides)} slides from {input_file} ---", file=sys.stderr)

if __name__ == "__main__":
    main()
